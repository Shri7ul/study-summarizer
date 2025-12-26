# PPTX reading logic
from pptx import Presentation


def read_pptx(file_path):
    slides_text = []

    prs = Presentation(file_path)

    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_content.append(text)

        if slide_content:
            slides_text.append(
                f"Slide {idx}:\n" + "\n".join(slide_content)
            )

    return slides_text
